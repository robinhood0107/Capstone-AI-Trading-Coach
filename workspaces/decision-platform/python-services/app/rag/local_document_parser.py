from __future__ import annotations

import hashlib
import io
import json
import math
import re
import stat
import zipfile
from dataclasses import dataclass
from html.parser import HTMLParser
from pathlib import Path, PurePosixPath
from typing import Any, Final, Literal, Protocol, cast

import fitz
from defusedxml import ElementTree as DefusedElementTree
from docx import Document
from docx.table import Table as DocxTable
from docx.text.paragraph import Paragraph as DocxParagraph
from openpyxl import load_workbook
from PIL import Image, UnidentifiedImageError
from pptx import Presentation

from app.rag.owner_file_io import OwnerFileIoError, read_owner_regular_file


class DocumentParseError(ValueError):
    """개인/OA 문서가 안전 처리·형식·자원 상한을 위반했을 때 stable code로 실패한다."""


BlockType = Literal["HEADING", "PARAGRAPH", "LIST", "TABLE", "FORMULA", "CAPTION"]

_PARSER_VERSION: Final = "1.1.0"
_PARSER_ARTIFACT_SHA256: Final = hashlib.sha256(
    b"capstone-s4-7d-safe-local-document-parser-v1.1"
).hexdigest()
_SOURCE_ID = re.compile(r"^src_[a-z0-9][a-z0-9_-]{2,95}$")
_REVISION_ID = re.compile(r"^srv_[a-z0-9][a-z0-9_-]{2,95}$")
_LANGUAGE_TAG = re.compile(r"^[a-z]{2,3}(?:-[A-Z]{2})?$")
_SUPPORTED_EXTENSIONS: Final[dict[str, str]] = {
    ".pdf": "application/pdf",
    ".docx": "application/vnd.openxmlformats-officedocument.wordprocessingml.document",
    ".pptx": "application/vnd.openxmlformats-officedocument.presentationml.presentation",
    ".xlsx": "application/vnd.openxmlformats-officedocument.spreadsheetml.sheet",
    ".html": "text/html",
    ".htm": "text/html",
    ".md": "text/markdown",
    ".markdown": "text/markdown",
    ".txt": "text/plain",
    ".png": "image/png",
    ".jpg": "image/jpeg",
    ".jpeg": "image/jpeg",
    ".tif": "image/tiff",
    ".tiff": "image/tiff",
}
_OOXML_MIME_MARKERS: Final[dict[str, bytes]] = {
    "application/vnd.openxmlformats-officedocument.wordprocessingml.document": (
        b"application/vnd.openxmlformats-officedocument.wordprocessingml.document.main+xml"
    ),
    "application/vnd.openxmlformats-officedocument.presentationml.presentation": (
        b"application/vnd.openxmlformats-officedocument.presentationml.presentation.main+xml"
    ),
    "application/vnd.openxmlformats-officedocument.spreadsheetml.sheet": (
        b"application/vnd.openxmlformats-officedocument.spreadsheetml.sheet.main+xml"
    ),
}
_PDF_ACTION_NAMES: Final[frozenset[str]] = frozenset(
    {
        "/GoTo",
        "/GoToR",
        "/GoToE",
        "/Launch",
        "/Thread",
        "/URI",
        "/Sound",
        "/Movie",
        "/Hide",
        "/Named",
        "/SubmitForm",
        "/ResetForm",
        "/ImportData",
        "/JavaScript",
        "/SetOCGState",
        "/Rendition",
        "/Trans",
        "/GoTo3DView",
    }
)
_PDF_EXECUTABLE_ACTION_NAMES: Final[frozenset[str]] = frozenset({"/JavaScript", "/Launch"})
_PDF_SAFE_OPEN_ACTION_NAMES: Final[frozenset[str]] = frozenset({"/GoTo"})
_PDF_XREF_REFERENCE: Final[re.Pattern[str]] = re.compile(r"^([1-9][0-9]*) 0 R$")
_PDF_OCR_RASTER_SCALE: Final[float] = 300 / 72
_MACRO_NAMES: Final[tuple[str, ...]] = (
    "vbaproject.bin",
    "vbadata.xml",
    "activex/",
    "embeddings/",
)
_SECRET_PATTERNS: Final[tuple[re.Pattern[str], ...]] = (
    re.compile(r"\b(?:api[_-]?key|secret|password)\s*[:=]\s*[\"']?[A-Za-z0-9_\-/.]{16,}", re.I),
    re.compile(r"\bsk-(?:proj-)?[A-Za-z0-9_-]{20,}\b"),
    re.compile(r"-----BEGIN (?:RSA |EC |OPENSSH )?PRIVATE KEY-----"),
    re.compile(r"\beyJ[A-Za-z0-9_-]{10,}\.[A-Za-z0-9_-]{10,}\.[A-Za-z0-9_-]{10,}\b"),
)
_PII_PATTERNS: Final[tuple[re.Pattern[str], ...]] = (
    re.compile(r"\b[A-Za-z0-9._%+-]+@[A-Za-z0-9.-]+\.[A-Za-z]{2,}\b"),
    re.compile(r"\b01[016789][ -]?\d{3,4}[ -]?\d{4}\b"),
    re.compile(r"\b\d{6}[ -]?[1-4]\d{6}\b"),
)
_PROMPT_INJECTION_PATTERNS: Final[tuple[re.Pattern[str], ...]] = (
    re.compile(r"ignore\s+(?:all\s+)?previous\s+instructions", re.I),
    re.compile(r"reveal\s+(?:the\s+)?system\s+prompt", re.I),
    re.compile(r"이전\s*(?:모든\s*)?지시(?:를|사항을)?\s*무시"),
    re.compile(r"<\|(?:system|assistant|user)\|>", re.I),
)


@dataclass(frozen=True, slots=True)
class ParserLimits:
    """문서 하나가 parser/OCR process에 소비할 수 있는 정적 상한이다."""

    max_file_bytes: int = 100 * 1024 * 1024
    max_archive_entries: int = 10_000
    max_decompressed_bytes: int = 256 * 1024 * 1024
    max_compression_ratio: int = 100
    max_pages: int = 500
    max_image_pixels: int = 100_000_000
    max_blocks: int = 50_000
    max_table_cells: int = 50_000
    max_text_characters: int = 10_000_000

    def validate(self) -> None:
        if any(
            value <= 0
            for value in (
                self.max_file_bytes,
                self.max_archive_entries,
                self.max_decompressed_bytes,
                self.max_compression_ratio,
                self.max_pages,
                self.max_image_pixels,
                self.max_blocks,
                self.max_table_cells,
                self.max_text_characters,
            )
        ):
            raise DocumentParseError("PARSER_LIMIT_INVALID")


@dataclass(slots=True)
class _TableCellBudget:
    """문서 전체의 table area를 native/OCR matrix 생성 전에 한 번만 선점한다."""

    max_cells: int
    consumed_cells: int = 0

    def reserve(self, *, row_count: object, column_count: object) -> None:
        """계산 중 overflow·큰 중간 allocation 없이 누적 row×column 예산을 닫는다."""

        if (
            type(row_count) is not int
            or type(column_count) is not int
            or row_count < 1
            or column_count < 1
            or row_count > self.max_cells
            or column_count > self.max_cells
            or row_count > (self.max_cells - self.consumed_cells) // column_count
        ):
            raise DocumentParseError("DOCUMENT_TABLE_BOUND_EXCEEDED")
        self.consumed_cells += row_count * column_count


@dataclass(slots=True)
class _BlockBudget:
    """Document IR block을 만들기 전에 문서 전체 개수 상한을 선점한다."""

    max_blocks: int
    consumed_blocks: int = 0

    def reserve(self) -> None:
        if self.consumed_blocks >= self.max_blocks:
            raise DocumentParseError("DOCUMENT_BLOCK_BOUND_EXCEEDED")
        self.consumed_blocks += 1


@dataclass(frozen=True, slots=True)
class OcrBlock:
    """OCR backend가 반환하는 경로·provider payload 없는 단일 구조 block이다."""

    block_type: BlockType
    confidence: float
    text: str | None = None
    level: int | None = None
    items: tuple[str, ...] = ()
    ordered: bool = False
    cells: tuple[tuple[int, int, str], ...] = ()
    row_count: int | None = None
    column_count: int | None = None
    normalized_formula: str | None = None
    target_reading_order: int | None = None


@dataclass(frozen=True, slots=True)
class OcrPageResult:
    """한 raster page의 OCR 결과이며 원본 image bytes나 backend body를 보존하지 않는다."""

    blocks: tuple[OcrBlock, ...]


class OcrBackendPort(Protocol):
    """network 없는 pinned production OCR process를 호출하는 내부 port다."""

    backend: str
    backend_version: str
    model_sha256: str

    def parse_page(self, *, png_bytes: bytes, page_number: int) -> OcrPageResult: ...


class LocalDocumentParser:
    """owner 파일을 복사·삭제하지 않고 native-first Document IR로 변환한다.

    입력 경로는 approved root 아래 descriptor/handle로만 읽고, 반환값에는 경로와 raw bytes를
    포함하지 않는다. 이미지 또는 text layer 없는 PDF page만 pinned OCR port로 전달한다.
    """

    def __init__(
        self,
        *,
        ocr_backend: OcrBackendPort | None = None,
        limits: ParserLimits | None = None,
        strip_inert_pdf_attachments: bool = False,
    ) -> None:
        """기본 owner 경계는 attachment를 거부하고, 승인된 OA 경계만 비실행 attachment를 제거한다."""

        if type(strip_inert_pdf_attachments) is not bool:
            raise DocumentParseError("PARSER_ATTACHMENT_POLICY_INVALID")
        self._ocr_backend = ocr_backend
        self._limits = limits or ParserLimits()
        self._strip_inert_pdf_attachments = strip_inert_pdf_attachments
        self._limits.validate()

    def parse_owner_document(
        self,
        *,
        approved_root: Path,
        relative_path: str,
        source_id: str,
        source_revision_id: str,
        language_tags: tuple[str, ...],
    ) -> dict[str, Any]:
        """owner가 보유한 regular file을 read-only로 파싱해 v1 Document IR dict를 반환한다."""

        return self.parse_approved_document(
            approved_root=approved_root,
            relative_path=relative_path,
            source_id=source_id,
            source_revision_id=source_revision_id,
            language_tags=language_tags,
        )

    def parse_approved_document(
        self,
        *,
        approved_root: Path,
        relative_path: str,
        source_id: str,
        source_revision_id: str,
        language_tags: tuple[str, ...],
    ) -> dict[str, Any]:
        """approved local regular file을 path-free Document IR로 변환한다.

        owner import와 OA raw cache는 같은 descriptor-safe regular-file reader를 사용한다. 이
        entrypoint는 source 종류에 관계없이 raw path·bytes를 결과에 넣지 않고, parser/OCR의
        resource 및 active-content 경계를 동일하게 적용한다.
        """

        _validate_identity(source_id, source_revision_id, language_tags)
        try:
            read_result = read_owner_regular_file(
                approved_root=approved_root,
                relative_path=relative_path,
                max_bytes=self._limits.max_file_bytes,
            )
        except OwnerFileIoError as error:
            raise DocumentParseError("DOCUMENT_PATH_UNSAFE") from error

        mime_type = _detect_mime(relative_path, read_result.content)
        archive: _ValidatedArchive | None = None
        if mime_type in _OOXML_MIME_MARKERS:
            archive = _validate_openxml_archive(
                read_result.content,
                mime_type=mime_type,
                limits=self._limits,
            )

        # 원본 parser가 block/table 객체를 만들기 전에 문서 전체 예산을 한 번만 공유한다.
        block_budget = _BlockBudget(max_blocks=self._limits.max_blocks)
        table_budget = _TableCellBudget(max_cells=self._limits.max_table_cells)
        blocks: list[dict[str, Any]]
        ocr_used = False
        if mime_type == "application/pdf":
            blocks, ocr_used = self._parse_pdf(
                read_result.content,
                block_budget=block_budget,
                table_budget=table_budget,
            )
        elif mime_type == "application/vnd.openxmlformats-officedocument.wordprocessingml.document":
            assert archive is not None
            blocks = _parse_docx(
                read_result.content,
                block_budget=block_budget,
                table_budget=table_budget,
            )
        elif mime_type == "application/vnd.openxmlformats-officedocument.presentationml.presentation":
            assert archive is not None
            blocks = _parse_pptx(
                read_result.content,
                block_budget=block_budget,
                table_budget=table_budget,
            )
        elif mime_type == "application/vnd.openxmlformats-officedocument.spreadsheetml.sheet":
            assert archive is not None
            blocks = _parse_xlsx(
                read_result.content,
                block_budget=block_budget,
                table_budget=table_budget,
            )
        elif mime_type == "text/html":
            blocks = _parse_html(
                _decode_utf8(read_result.content),
                block_budget=block_budget,
                table_budget=table_budget,
            )
        elif mime_type == "text/markdown":
            blocks = _parse_markdown(
                _decode_utf8(read_result.content),
                block_budget=block_budget,
                table_budget=table_budget,
            )
        elif mime_type == "text/plain":
            blocks = _parse_plain_text(
                _decode_utf8(read_result.content),
                block_budget=block_budget,
            )
        elif mime_type.startswith("image/"):
            blocks = self._parse_image(
                read_result.content,
                block_budget=block_budget,
                table_budget=table_budget,
            )
            ocr_used = True
        else:  # pragma: no cover - closed MIME set is enforced above.
            raise DocumentParseError("DOCUMENT_MIME_UNSUPPORTED")

        _validate_block_bounds(blocks, self._limits)
        safety = _classify_safety(blocks)
        if safety["secretDetected"]:
            raise DocumentParseError("DOCUMENT_SECRET_QUARANTINED")
        normalized = json.dumps(blocks, ensure_ascii=False, separators=(",", ":"), sort_keys=True)
        extraction_mode = _extraction_mode(blocks, ocr_used)
        ocr_evidence = self._ocr_evidence(ocr_used)
        return {
            "blocks": blocks,
            "contractId": "rag-document-ir-v1",
            "documentIrVersion": 1,
            "extractionMode": extraction_mode,
            "languageTags": list(language_tags),
            "mimeType": mime_type,
            "normalizedContentSha256": hashlib.sha256(normalized.encode("utf-8")).hexdigest(),
            "parserEvidence": {
                "ocr": ocr_evidence,
                "parserArtifactSha256": _PARSER_ARTIFACT_SHA256,
                "parserBackend": "capstone-safe-local-document-parser",
                "parserVersion": _PARSER_VERSION,
            },
            "rawContentSha256": read_result.content_sha256,
            "safetyClassification": safety,
            "sourceId": source_id,
            "sourceRevisionId": source_revision_id,
        }

    def _parse_pdf(
        self,
        payload: bytes,
        *,
        block_budget: _BlockBudget,
        table_budget: _TableCellBudget,
    ) -> tuple[list[dict[str, Any]], bool]:
        document = _open_normalized_pdf(payload, self._limits)
        try:
            has_inert_attachments = _inspect_pdf_objects(
                document,
                self._limits,
                allow_inert_attachments=self._strip_inert_pdf_attachments,
            )
            if has_inert_attachments:
                expected_page_count = document.page_count
                sanitized_payload = _strip_inert_pdf_attachments(document, self._limits)
                document.close()
                document = _open_normalized_pdf(sanitized_payload, self._limits)
                if document.page_count != expected_page_count or _inspect_pdf_objects(
                    document,
                    self._limits,
                    allow_inert_attachments=False,
                ):
                    raise DocumentParseError("PDF_ATTACHMENT_SANITIZATION_FAILED")
            output: list[dict[str, Any]] = []
            ocr_used = False
            for page_number, page in enumerate(document, start=1):
                native_blocks = _native_pdf_blocks(
                    page,
                    page_number,
                    block_budget=block_budget,
                )
                if native_blocks:
                    output.extend(native_blocks)
                    continue
                output.extend(
                    self._ocr_page(
                        _render_pdf_page(page, max_image_pixels=self._limits.max_image_pixels),
                        page_number,
                        {"page": page_number},
                        block_budget=block_budget,
                        table_budget=table_budget,
                    )
                )
                ocr_used = True
            return _renumber(output), ocr_used
        finally:
            document.close()

    def _parse_image(
        self,
        payload: bytes,
        *,
        block_budget: _BlockBudget,
        table_budget: _TableCellBudget,
    ) -> list[dict[str, Any]]:
        try:
            image = Image.open(io.BytesIO(payload))
        except (UnidentifiedImageError, OSError) as error:
            raise DocumentParseError("IMAGE_MALFORMED") from error
        output: list[dict[str, Any]] = []
        try:
            frames = getattr(image, "n_frames", 1)
            if frames < 1 or frames > self._limits.max_pages:
                raise DocumentParseError("IMAGE_PAGE_BOUND_EXCEEDED")
            for page_number in range(1, frames + 1):
                image.seek(page_number - 1)
                if image.width * image.height > self._limits.max_image_pixels:
                    raise DocumentParseError("IMAGE_PIXEL_BOUND_EXCEEDED")
                converted = io.BytesIO()
                image.convert("RGB").save(converted, format="PNG")
                output.extend(
                    self._ocr_page(
                        converted.getvalue(),
                        page_number,
                        {"page": page_number},
                        block_budget=block_budget,
                        table_budget=table_budget,
                    )
                )
        finally:
            image.close()
        return _renumber(output)

    def _ocr_page(
        self,
        png_bytes: bytes,
        page_number: int,
        locator: dict[str, object],
        *,
        block_budget: _BlockBudget,
        table_budget: _TableCellBudget,
    ) -> list[dict[str, Any]]:
        if self._ocr_backend is None:
            raise DocumentParseError("OCR_BACKEND_REQUIRED")
        _validate_ocr_identity(self._ocr_backend)
        try:
            result = self._ocr_backend.parse_page(
                png_bytes=png_bytes,
                page_number=page_number,
            )
        except DocumentParseError:
            raise
        except Exception as error:
            raise DocumentParseError("OCR_BACKEND_FAILED") from error
        if not isinstance(result, OcrPageResult) or not result.blocks:
            raise DocumentParseError("OCR_EMPTY_RESULT")
        return [
            _ocr_block_to_ir(
                block,
                locator,
                block_budget=block_budget,
                table_budget=table_budget,
            )
            for block in result.blocks
        ]

    def _ocr_evidence(self, ocr_used: bool) -> dict[str, object]:
        if not ocr_used:
            return {"backend": "NOT_USED", "backendVersion": None, "modelSha256": None}
        if self._ocr_backend is None:  # pragma: no cover - guarded by _ocr_page.
            raise DocumentParseError("OCR_BACKEND_REQUIRED")
        _validate_ocr_identity(self._ocr_backend)
        return {
            "backend": self._ocr_backend.backend,
            "backendVersion": self._ocr_backend.backend_version,
            "modelSha256": self._ocr_backend.model_sha256,
        }


@dataclass(frozen=True, slots=True)
class _ValidatedArchive:
    names: tuple[str, ...]


def _validate_identity(
    source_id: str,
    source_revision_id: str,
    language_tags: tuple[str, ...],
) -> None:
    if (
        _SOURCE_ID.fullmatch(source_id) is None
        or _REVISION_ID.fullmatch(source_revision_id) is None
        or not 1 <= len(language_tags) <= 10
        or len(set(language_tags)) != len(language_tags)
        or any(_LANGUAGE_TAG.fullmatch(tag) is None for tag in language_tags)
    ):
        raise DocumentParseError("DOCUMENT_IDENTITY_INVALID")


def _detect_mime(relative_path: str, payload: bytes) -> str:
    suffix = PurePosixPath(relative_path).suffix.lower()
    mime = _SUPPORTED_EXTENSIONS.get(suffix)
    if mime is None:
        raise DocumentParseError("DOCUMENT_MIME_UNSUPPORTED")
    if mime == "application/pdf":
        valid = payload.startswith(b"%PDF-") and b"%%EOF" in payload[-4096:]
    elif mime in _OOXML_MIME_MARKERS:
        valid = payload.startswith(b"PK\x03\x04")
    elif mime == "image/png":
        valid = payload.startswith(b"\x89PNG\r\n\x1a\n")
    elif mime == "image/jpeg":
        valid = payload.startswith(b"\xff\xd8\xff") and payload.endswith(b"\xff\xd9")
    elif mime == "image/tiff":
        valid = payload.startswith((b"II*\x00", b"MM\x00*"))
    else:
        try:
            text = _decode_utf8(payload)
        except DocumentParseError:
            valid = False
        else:
            lowered = text.lstrip().lower()
            valid = mime != "text/html" or lowered.startswith(("<!doctype html", "<html", "<head", "<body"))
    if not valid:
        raise DocumentParseError("DOCUMENT_MIME_MISMATCH")
    return mime


def _validate_openxml_archive(
    payload: bytes,
    *,
    mime_type: str,
    limits: ParserLimits,
) -> _ValidatedArchive:
    try:
        archive = zipfile.ZipFile(io.BytesIO(payload))
    except zipfile.BadZipFile as error:
        raise DocumentParseError("ARCHIVE_MALFORMED") from error
    try:
        infos = archive.infolist()
        if not 1 <= len(infos) <= limits.max_archive_entries:
            raise DocumentParseError("ARCHIVE_ENTRY_BOUND_EXCEEDED")
        names: set[str] = set()
        total = 0
        for info in infos:
            name = info.filename
            _validate_archive_name(name)
            folded = name.casefold()
            if folded in names:
                raise DocumentParseError("ARCHIVE_DUPLICATE_MEMBER")
            names.add(folded)
            mode = info.external_attr >> 16
            if stat.S_ISLNK(mode):
                raise DocumentParseError("ARCHIVE_SYMLINK_FORBIDDEN")
            if info.flag_bits & 0x1:
                raise DocumentParseError("ARCHIVE_ENCRYPTION_FORBIDDEN")
            total += info.file_size
            if total > limits.max_decompressed_bytes:
                raise DocumentParseError("ARCHIVE_DECOMPRESSED_BOUND_EXCEEDED")
            if info.file_size and info.file_size / max(1, info.compress_size) > limits.max_compression_ratio:
                raise DocumentParseError("ARCHIVE_COMPRESSION_RATIO_EXCEEDED")
            if any(token in folded for token in _MACRO_NAMES):
                raise DocumentParseError("OFFICE_MACRO_FORBIDDEN")
            if "externallinks/" in folded:
                raise DocumentParseError("OFFICE_EXTERNAL_RELATIONSHIP_FORBIDDEN")
            if folded.endswith((".xml", ".rels")):
                content = archive.read(info)
                lowered = content.lower()
                if b"<!doctype" in lowered or b"<!entity" in lowered:
                    raise DocumentParseError("XML_DTD_FORBIDDEN")
                if folded.endswith(".rels"):
                    _validate_relationships(content)
        try:
            content_types = archive.read("[Content_Types].xml")
        except KeyError as error:
            raise DocumentParseError("OPENXML_CONTENT_TYPE_MISSING") from error
        if _OOXML_MIME_MARKERS[mime_type] not in content_types:
            raise DocumentParseError("DOCUMENT_MIME_MISMATCH")
        return _ValidatedArchive(names=tuple(sorted(names)))
    except (OSError, RuntimeError, zipfile.BadZipFile) as error:
        raise DocumentParseError("ARCHIVE_MALFORMED") from error
    finally:
        archive.close()


def _validate_archive_name(name: str) -> None:
    if (
        not name
        or "\\" in name
        or name.startswith("/")
        or "\x00" in name
        or ":" in name.split("/", 1)[0]
        or any(part in {"", ".", ".."} for part in name.rstrip("/").split("/"))
    ):
        raise DocumentParseError("ARCHIVE_PATH_UNSAFE")


def _validate_relationships(payload: bytes) -> None:
    try:
        root = DefusedElementTree.fromstring(payload)
    except Exception as error:
        raise DocumentParseError("XML_MALFORMED") from error
    for relationship in root.iter():
        mode = relationship.attrib.get("TargetMode", "")
        target = relationship.attrib.get("Target", "")
        if mode.casefold() == "external" or target.lower().startswith(
            ("http:", "https:", "file:", "ftp:", "mailto:", "\\\\")
        ):
            raise DocumentParseError("OFFICE_EXTERNAL_RELATIONSHIP_FORBIDDEN")


def _decode_utf8(payload: bytes) -> str:
    try:
        value = payload.decode("utf-8-sig")
    except UnicodeDecodeError as error:
        raise DocumentParseError("TEXT_ENCODING_UNSUPPORTED") from error
    if "\x00" in value:
        raise DocumentParseError("TEXT_NUL_FORBIDDEN")
    return value.replace("\r\n", "\n").replace("\r", "\n")


def _parse_plain_text(text: str, *, block_budget: _BlockBudget) -> list[dict[str, Any]]:
    blocks: list[dict[str, Any]] = []
    start = 0
    for boundary in re.finditer(r"\n\s*\n", text):
        paragraph = text[start : boundary.start()].strip()
        if paragraph:
            block_budget.reserve()
            blocks.append(_paragraph({"section": "document"}, paragraph))
        start = boundary.end()
    paragraph = text[start:].strip()
    if paragraph:
        block_budget.reserve()
        blocks.append(_paragraph({"section": "document"}, paragraph))
    if not blocks:
        raise DocumentParseError("DOCUMENT_EMPTY")
    return _renumber(blocks)


def _parse_markdown(
    text: str,
    *,
    block_budget: _BlockBudget,
    table_budget: _TableCellBudget,
) -> list[dict[str, Any]]:
    output: list[dict[str, Any]] = []
    section = "document"
    paragraph: list[str] = []
    list_items: list[str] = []
    table_lines: list[str] = []

    def flush_paragraph() -> None:
        if paragraph:
            block_budget.reserve()
            output.append(_paragraph({"section": section}, "\n".join(paragraph)))
            paragraph.clear()

    def flush_list() -> None:
        if list_items:
            block_budget.reserve()
            output.append(_list_block({"section": section}, tuple(list_items), ordered=False))
            list_items.clear()

    def flush_table() -> None:
        if table_lines:
            block_budget.reserve()
            output.append(
                _markdown_table(
                    {"section": section},
                    table_lines,
                    table_budget=table_budget,
                )
            )
            table_lines.clear()

    for raw_line in text.splitlines():
        line = raw_line.rstrip()
        heading = re.match(r"^(#{1,6})\s+(.+?)\s*$", line)
        if heading:
            flush_paragraph()
            flush_list()
            flush_table()
            section = heading.group(2).strip()
            block_budget.reserve()
            output.append(_heading({"section": section}, section, len(heading.group(1))))
        elif re.match(r"^\s*[-*+]\s+\S", line):
            flush_paragraph()
            flush_table()
            list_items.append(re.sub(r"^\s*[-*+]\s+", "", line).strip())
        elif line.strip().startswith("|") and line.strip().endswith("|"):
            flush_paragraph()
            flush_list()
            table_lines.append(line.strip())
        elif not line.strip():
            flush_paragraph()
            flush_list()
            flush_table()
        else:
            flush_list()
            flush_table()
            paragraph.append(line)
    flush_paragraph()
    flush_list()
    flush_table()
    if not output:
        raise DocumentParseError("DOCUMENT_EMPTY")
    return _renumber(output)


class _SafeHtmlParser(HTMLParser):
    def __init__(
        self,
        *,
        block_budget: _BlockBudget,
        table_budget: _TableCellBudget,
    ) -> None:
        super().__init__(convert_charrefs=True)
        self._block_budget = block_budget
        self._table_budget = table_budget
        self.blocks: list[dict[str, Any]] = []
        self._tag: str | None = None
        self._buffer: list[str] = []
        self._list_items: list[str] = []
        self._table_rows: list[list[str]] = []
        self._row: list[str] = []
        self._section = "document"

    def handle_starttag(self, tag: str, attrs: list[tuple[str, str | None]]) -> None:
        folded = tag.casefold()
        attributes = {key.casefold(): value or "" for key, value in attrs}
        if (
            folded in {"script", "iframe", "object", "embed", "base", "link"}
            or any(key.startswith("on") for key in attributes)
            or (folded in {"img", "audio", "video", "source", "form"} and any(
                key in attributes for key in ("src", "srcset", "action", "poster")
            ))
        ):
            raise DocumentParseError("HTML_ACTIVE_RESOURCE_FORBIDDEN")
        if folded in {"h1", "h2", "h3", "h4", "h5", "h6", "p", "li", "td", "th", "caption"}:
            self._tag = folded
            self._buffer = []
        if folded == "tr":
            self._row = []

    def handle_endtag(self, tag: str) -> None:
        folded = tag.casefold()
        text = " ".join("".join(self._buffer).split())
        if folded.startswith("h") and len(folded) == 2 and folded[1].isdigit() and text:
            self._section = text
            self._block_budget.reserve()
            self.blocks.append(_heading({"section": text}, text, int(folded[1])))
        elif folded == "p" and text:
            self._block_budget.reserve()
            self.blocks.append(_paragraph({"section": self._section}, text))
        elif folded == "li" and text:
            self._list_items.append(text)
        elif folded in {"td", "th"} and text:
            self._row.append(text)
        elif folded == "tr" and self._row:
            self._table_rows.append(self._row)
            self._row = []
        elif folded in {"ul", "ol"} and self._list_items:
            self._block_budget.reserve()
            self.blocks.append(
                _list_block(
                    {"section": self._section},
                    tuple(self._list_items),
                    ordered=folded == "ol",
                )
            )
            self._list_items = []
        elif folded == "table" and self._table_rows:
            self._block_budget.reserve()
            self.blocks.append(
                _table(
                    {"section": self._section},
                    self._table_rows,
                    table_budget=self._table_budget,
                )
            )
            self._table_rows = []
        self._tag = None
        self._buffer = []

    def handle_data(self, data: str) -> None:
        if self._tag == "style" and re.search(r"url\s*\(|@import", data, re.I):
            raise DocumentParseError("HTML_ACTIVE_RESOURCE_FORBIDDEN")
        if self._tag is not None:
            self._buffer.append(data)


def _parse_html(
    text: str,
    *,
    block_budget: _BlockBudget,
    table_budget: _TableCellBudget,
) -> list[dict[str, Any]]:
    if re.search(r"<(?:script|iframe|object|embed|link|base)\b", text, re.I) or re.search(
        r"(?:url\s*\(|@import)", text, re.I
    ):
        raise DocumentParseError("HTML_ACTIVE_RESOURCE_FORBIDDEN")
    parser = _SafeHtmlParser(block_budget=block_budget, table_budget=table_budget)
    try:
        parser.feed(text)
        parser.close()
    except DocumentParseError:
        raise
    except Exception as error:
        raise DocumentParseError("HTML_MALFORMED") from error
    if not parser.blocks:
        raise DocumentParseError("DOCUMENT_EMPTY")
    return _renumber(parser.blocks)


def _open_normalized_pdf(payload: bytes, limits: ParserLimits) -> fitz.Document:
    """reachable object만 메모리에서 재직렬화해 과거 incremental xref를 제거한다."""

    normalized: fitz.Document | None = None
    try:
        source = fitz.open(stream=payload, filetype="pdf")
    except Exception as error:
        raise DocumentParseError("PDF_MALFORMED") from error
    try:
        _validate_pdf_container(source, limits)
        normalized_payload = _serialize_normalized_pdf(source, limits)
    finally:
        source.close()
    try:
        normalized = fitz.open(stream=normalized_payload, filetype="pdf")
        _validate_pdf_container(normalized, limits)
    except DocumentParseError:
        if normalized is not None:
            normalized.close()
        raise
    except Exception as error:
        if normalized is not None:
            normalized.close()
        raise DocumentParseError("PDF_MALFORMED") from error
    return normalized


def _validate_pdf_container(document: fitz.Document, limits: ParserLimits) -> None:
    """정규화 전후에 encryption, page, object 상한을 동일하게 적용한다."""

    try:
        if document.needs_pass or document.is_encrypted:
            raise DocumentParseError("PDF_ENCRYPTED_FORBIDDEN")
        object_count = document.xref_length()
        if object_count < 1 or object_count > limits.max_archive_entries:
            raise DocumentParseError("PDF_OBJECT_BOUND_EXCEEDED")
        if document.page_count < 1 or document.page_count > limits.max_pages:
            raise DocumentParseError("PDF_PAGE_BOUND_EXCEEDED")
    except DocumentParseError:
        raise
    except Exception as error:
        raise DocumentParseError("PDF_MALFORMED") from error


def _serialize_normalized_pdf(document: fitz.Document, limits: ParserLimits) -> bytes:
    """raw cache를 바꾸지 않고 reachable PDF projection만 bounded memory bytes로 만든다."""

    try:
        payload = cast(
            bytes,
            document.tobytes(
                garbage=4,
                clean=True,
                deflate=True,
                deflate_images=True,
                deflate_fonts=True,
            ),
        )
    except Exception as error:
        raise DocumentParseError("PDF_MALFORMED") from error
    if not payload or len(payload) > limits.max_decompressed_bytes:
        raise DocumentParseError("PDF_NORMALIZED_BOUND_EXCEEDED")
    return payload


def _inspect_pdf_objects(
    document: fitz.Document,
    limits: ParserLimits,
    *,
    allow_inert_attachments: bool,
) -> bool:
    """구조 key를 검사해 실행 action을 거부하고 inert attachment 존재 여부만 반환한다."""

    _validate_pdf_container(document, limits)
    try:
        catalog_xref = document.pdf_catalog()
        embedded_file_count = document.embfile_count()
    except Exception as error:
        raise DocumentParseError("PDF_MALFORMED") from error
    if _pdf_key(document, catalog_xref, "AA")[0] != "null":
        raise DocumentParseError("PDF_ACTIVE_CONTENT_FORBIDDEN")
    if _pdf_key(document, catalog_xref, "AcroForm/XFA")[0] != "null":
        raise DocumentParseError("PDF_ACTIVE_CONTENT_FORBIDDEN")
    _validate_pdf_open_action(document, catalog_xref)

    attachment_object_seen = False
    object_count = document.xref_length()
    for xref in range(1, object_count):
        action = _pdf_key(document, xref, "S")
        nested_action = _pdf_key(document, xref, "A")
        javascript = _pdf_key(document, xref, "JS")
        additional_action = _pdf_key(document, xref, "AA")
        object_type = _pdf_key(document, xref, "Type")
        subtype = _pdf_key(document, xref, "Subtype")
        open_action = _pdf_key(document, xref, "OpenAction")
        xfa = _pdf_key(document, xref, "XFA")
        embedded_file = _pdf_key(document, xref, "EF")
        associated_file = _pdf_key(document, xref, "AF")
        if (
            (action[0] == "name" and action[1] in _PDF_EXECUTABLE_ACTION_NAMES)
            or _pdf_nested_action_is_forbidden(document, xref, nested_action)
            or (
                action[0] == "name"
                and action[1] in _PDF_ACTION_NAMES
                and _pdf_key(document, xref, "Next")[0] != "null"
            )
            or javascript[0] != "null"
            or additional_action[0] != "null"
            or object_type == ("name", "/RichMedia")
            or subtype == ("name", "/RichMedia")
            or (xref != catalog_xref and open_action[0] != "null")
            or xfa[0] != "null"
        ):
            raise DocumentParseError("PDF_ACTIVE_CONTENT_FORBIDDEN")
        is_file_spec = object_type == ("name", "/Filespec")
        is_embedded_stream = object_type == ("name", "/EmbeddedFile")
        is_file_attachment = subtype == ("name", "/FileAttachment")
        if is_file_spec or is_embedded_stream or is_file_attachment:
            # 외부 file spec은 inert attachment가 아니므로 OA 경계에서도 허용하지 않는다.
            if is_file_spec and embedded_file[0] == "null":
                raise DocumentParseError("PDF_ATTACHMENT_FORBIDDEN")
            attachment_object_seen = True
        if associated_file[0] != "null":
            attachment_object_seen = True

    if _pdf_name_tree_populated(document, catalog_xref, "JavaScript"):
        raise DocumentParseError("PDF_ACTIVE_CONTENT_FORBIDDEN")
    if _pdf_name_tree_populated(document, catalog_xref, "EmbeddedFiles"):
        attachment_object_seen = True
    has_inert_attachments = embedded_file_count > 0 or attachment_object_seen
    if has_inert_attachments and not allow_inert_attachments:
        raise DocumentParseError("PDF_ATTACHMENT_FORBIDDEN")
    return has_inert_attachments


def _pdf_nested_action_is_forbidden(
    document: fitz.Document,
    xref: int,
    action: tuple[str, str],
) -> bool:
    """annotation·outline의 direct action과 action chain에서 실행 payload를 닫는다."""

    if action[0] == "name":
        return action[1] in _PDF_EXECUTABLE_ACTION_NAMES
    if action[0] != "dict":
        return False
    nested_action = _pdf_key(document, xref, "A/S")
    return (
        (nested_action[0] == "name" and nested_action[1] in _PDF_EXECUTABLE_ACTION_NAMES)
        or _pdf_key(document, xref, "A/JS")[0] != "null"
        # `/Next`는 dictionary 또는 array action chain이라 재귀 우회를 막기 위해 전부 거부한다.
        or _pdf_key(document, xref, "A/Next")[0] != "null"
    )


def _validate_pdf_open_action(document: fitz.Document, catalog_xref: int) -> None:
    """initial page 이동만 허용하고 executable open action은 거부한다."""

    kind, value = _pdf_key(document, catalog_xref, "OpenAction")
    if kind == "null" or kind == "array":
        return
    if kind != "xref":
        raise DocumentParseError("PDF_ACTIVE_CONTENT_FORBIDDEN")
    match = _PDF_XREF_REFERENCE.fullmatch(value)
    if match is None:
        raise DocumentParseError("PDF_ACTIVE_CONTENT_FORBIDDEN")
    target_xref = int(match.group(1))
    action = _pdf_key(document, target_xref, "S")
    destination = _pdf_key(document, target_xref, "D")
    if (
        action[0] != "name"
        or action[1] not in _PDF_SAFE_OPEN_ACTION_NAMES
        or destination[0] not in {"array", "name", "string", "xref"}
        or _pdf_key(document, target_xref, "JS")[0] != "null"
        or _pdf_key(document, target_xref, "AA")[0] != "null"
    ):
        raise DocumentParseError("PDF_ACTIVE_CONTENT_FORBIDDEN")


def _strip_inert_pdf_attachments(document: fitz.Document, limits: ParserLimits) -> bytes:
    """OA raw는 보존하고 in-memory projection에서 file payload와 annotation만 제거한다."""

    try:
        document.scrub(
            attached_files=True,
            clean_pages=False,
            embedded_files=True,
            hidden_text=False,
            javascript=True,
            metadata=False,
            redactions=False,
            remove_links=False,
            reset_fields=False,
            reset_responses=False,
            thumbnails=False,
            xml_metadata=False,
        )
        # PDF 2.0 수식 접근성 자료처럼 StructElem의 /AF로 연결된 파일은 scrub 대상이 아니다.
        # OA projection에서 참조를 먼저 끊어야 garbage=4가 Filespec과 EmbeddedFile을 제거한다.
        for xref in range(1, document.xref_length()):
            if _pdf_key(document, xref, "AF")[0] != "null":
                document.xref_set_key(xref, "AF", "null")
    except Exception as error:
        raise DocumentParseError("PDF_ATTACHMENT_SANITIZATION_FAILED") from error
    return _serialize_normalized_pdf(document, limits)


def _pdf_name_tree_populated(document: fitz.Document, catalog_xref: int, name: str) -> bool:
    """scrub가 남기는 empty `/Names []` dictionary는 active payload로 세지 않는다."""

    kind, _value = _pdf_key(document, catalog_xref, f"Names/{name}")
    if kind == "null":
        return False
    names_kind, names_value = _pdf_key(document, catalog_xref, f"Names/{name}/Names")
    return names_kind != "array" or names_value != "[]"


def _pdf_key(document: fitz.Document, xref: int, key: str) -> tuple[str, str]:
    """xref key read 실패를 malformed PDF stable code로 닫는다."""

    try:
        kind, value = document.xref_get_key(xref, key)
    except Exception as error:
        raise DocumentParseError("PDF_MALFORMED") from error
    if not isinstance(kind, str) or not isinstance(value, str):
        raise DocumentParseError("PDF_MALFORMED")
    return kind, value


def _native_pdf_blocks(
    page: fitz.Page,
    page_number: int,
    *,
    block_budget: _BlockBudget,
) -> list[dict[str, Any]]:
    output: list[dict[str, Any]] = []
    try:
        native = page.get_text("blocks", sort=True)
    except Exception as error:
        raise DocumentParseError("PDF_PARSE_FAILED") from error
    for item in native:
        text = " ".join(str(item[4]).split())
        if text:
            block_budget.reserve()
            output.append(_paragraph({"page": page_number}, text))
    return output


def _render_pdf_page(page: fitz.Page, *, max_image_pixels: int) -> bytes:
    try:
        _validate_pdf_raster_bounds(page, max_image_pixels=max_image_pixels)
        pixmap = page.get_pixmap(matrix=fitz.Matrix(_PDF_OCR_RASTER_SCALE, _PDF_OCR_RASTER_SCALE), alpha=False)
        if pixmap.width * pixmap.height > max_image_pixels:
            raise DocumentParseError("IMAGE_PIXEL_BOUND_EXCEEDED")
        return cast(bytes, pixmap.tobytes("png"))
    except DocumentParseError:
        raise
    except Exception as error:
        raise DocumentParseError("PDF_RENDER_FAILED") from error


def _validate_pdf_raster_bounds(page: fitz.Page, *, max_image_pixels: int) -> None:
    """PDF page geometry를 raster allocation 전에 검사해 OCR worker의 메모리 상한을 보존한다."""

    try:
        width_points = float(page.rect.width)
        height_points = float(page.rect.height)
    except (AttributeError, TypeError, ValueError, OverflowError) as error:
        raise DocumentParseError("PDF_RENDER_FAILED") from error
    if (
        not math.isfinite(width_points)
        or not math.isfinite(height_points)
        or width_points <= 0
        or height_points <= 0
    ):
        raise DocumentParseError("PDF_RENDER_FAILED")
    scaled_width = width_points * _PDF_OCR_RASTER_SCALE
    scaled_height = height_points * _PDF_OCR_RASTER_SCALE
    if not math.isfinite(scaled_width) or not math.isfinite(scaled_height):
        raise DocumentParseError("IMAGE_PIXEL_BOUND_EXCEEDED")
    width_pixels = math.ceil(scaled_width)
    height_pixels = math.ceil(scaled_height)
    if (
        width_pixels > max_image_pixels
        or height_pixels > max_image_pixels
        or height_pixels > max_image_pixels // width_pixels
    ):
        raise DocumentParseError("IMAGE_PIXEL_BOUND_EXCEEDED")


def _parse_docx(
    payload: bytes,
    *,
    block_budget: _BlockBudget,
    table_budget: _TableCellBudget,
) -> list[dict[str, Any]]:
    try:
        document = Document(io.BytesIO(payload))
    except Exception as error:
        raise DocumentParseError("DOCX_PARSE_FAILED") from error
    output: list[dict[str, Any]] = []
    section = "document"
    for child in document.element.body.iterchildren():
        if child.tag.endswith("}p"):
            paragraph = DocxParagraph(child, document)
            text = " ".join(paragraph.text.split())
            if not text:
                continue
            style = paragraph.style.name if paragraph.style is not None else ""
            match = re.fullmatch(r"Heading\s+([1-6])", style, re.I)
            block_budget.reserve()
            if match:
                section = text
                output.append(_heading({"section": section}, text, int(match.group(1))))
            elif paragraph.style is not None and "List" in style:
                output.append(_list_block({"section": section}, (text,), ordered="Number" in style))
            else:
                output.append(_paragraph({"section": section}, text))
        elif child.tag.endswith("}tbl"):
            block_budget.reserve()
            table = DocxTable(child, document)
            table_budget.reserve(
                row_count=len(table.rows),
                column_count=len(table.columns),
            )
            rows = [[" ".join(cell.text.split()) for cell in row.cells] for row in table.rows]
            output.append(
                _table(
                    {"section": section},
                    rows,
                    table_budget=table_budget,
                    pre_reserved=True,
                )
            )
    if not output:
        raise DocumentParseError("DOCUMENT_EMPTY")
    return _renumber(output)


def _parse_pptx(
    payload: bytes,
    *,
    block_budget: _BlockBudget,
    table_budget: _TableCellBudget,
) -> list[dict[str, Any]]:
    try:
        presentation = Presentation(io.BytesIO(payload))
    except Exception as error:
        raise DocumentParseError("PPTX_PARSE_FAILED") from error
    output: list[dict[str, Any]] = []
    for slide_number, slide in enumerate(presentation.slides, start=1):
        title_shape = slide.shapes.title
        for shape in slide.shapes:
            if getattr(shape, "has_table", False):
                block_budget.reserve()
                table_budget.reserve(
                    row_count=len(shape.table.rows),
                    column_count=len(shape.table.columns),
                )
                rows = [
                    [" ".join(cell.text.split()) for cell in row.cells]
                    for row in shape.table.rows
                ]
                output.append(
                    _table(
                        {"slide": slide_number},
                        rows,
                        table_budget=table_budget,
                        pre_reserved=True,
                    )
                )
                continue
            if not getattr(shape, "has_text_frame", False):
                continue
            text = " ".join(shape.text.split())
            if not text:
                continue
            block_budget.reserve()
            if shape == title_shape:
                output.append(_heading({"slide": slide_number}, text, 1))
            else:
                output.append(_paragraph({"slide": slide_number}, text))
    if not output:
        raise DocumentParseError("DOCUMENT_EMPTY")
    return _renumber(output)


def _parse_xlsx(
    payload: bytes,
    *,
    block_budget: _BlockBudget,
    table_budget: _TableCellBudget,
) -> list[dict[str, Any]]:
    try:
        workbook = load_workbook(
            io.BytesIO(payload),
            read_only=True,
            data_only=False,
            keep_links=False,
        )
    except Exception as error:
        raise DocumentParseError("XLSX_PARSE_FAILED") from error
    output: list[dict[str, Any]] = []
    try:
        for sheet in workbook.worksheets:
            # read_only iterator가 Python rows를 보유하기 전에 worksheet의 declared area를 닫는다.
            table_budget.reserve(
                row_count=sheet.max_row,
                column_count=sheet.max_column,
            )
            rows: list[list[str]] = []
            formulas: list[str] = []
            table_block_reserved = False
            for row in sheet.iter_rows(values_only=True):
                values: list[str] = []
                for value in row:
                    if value is None:
                        values.append("")
                    elif isinstance(value, str) and value.startswith("="):
                        block_budget.reserve()
                        formulas.append(value)
                        values.append(value)
                    else:
                        values.append(str(value))
                if any(values):
                    if not table_block_reserved:
                        block_budget.reserve()
                        table_block_reserved = True
                    rows.append(values)
            if rows:
                output.append(
                    _table(
                        {"sheet": sheet.title},
                        rows,
                        table_budget=table_budget,
                        pre_reserved=True,
                    )
                )
            for formula in formulas:
                output.append(_formula({"sheet": sheet.title}, formula, _normalize_formula(formula)))
    finally:
        workbook.close()
    if not output:
        raise DocumentParseError("DOCUMENT_EMPTY")
    return _renumber(output)


def _paragraph(locator: dict[str, object], text: str, confidence: float | None = None) -> dict[str, Any]:
    return {
        "blockType": "PARAGRAPH",
        "locator": locator,
        "ocrConfidence": confidence,
        "readingOrder": 0,
        "text": text.strip(),
    }


def _heading(
    locator: dict[str, object],
    text: str,
    level: int,
    confidence: float | None = None,
) -> dict[str, Any]:
    return {
        "blockType": "HEADING",
        "level": max(1, min(6, level)),
        "locator": locator,
        "ocrConfidence": confidence,
        "readingOrder": 0,
        "text": text.strip(),
    }


def _list_block(
    locator: dict[str, object],
    items: tuple[str, ...],
    *,
    ordered: bool,
    confidence: float | None = None,
) -> dict[str, Any]:
    return {
        "blockType": "LIST",
        "items": [item.strip() for item in items if item.strip()],
        "locator": locator,
        "ocrConfidence": confidence,
        "ordered": ordered,
        "readingOrder": 0,
    }


def _table(
    locator: dict[str, object],
    rows: list[list[str]],
    confidence: float | None = None,
    *,
    table_budget: _TableCellBudget,
    pre_reserved: bool = False,
) -> dict[str, Any]:
    if not rows:
        raise DocumentParseError("TABLE_EMPTY")
    source_column_count = max((len(row) for row in rows), default=0)
    if source_column_count < 1:
        raise DocumentParseError("TABLE_EMPTY")
    if not pre_reserved:
        table_budget.reserve(
            row_count=len(rows),
            column_count=source_column_count,
        )
    normalized_rows = [[value.strip() for value in row] for row in rows if any(value.strip() for value in row)]
    if not normalized_rows:
        raise DocumentParseError("TABLE_EMPTY")
    column_count = max(len(row) for row in normalized_rows)
    cells = [
        {
            "column": column_index,
            "columnSpan": 1,
            "row": row_index,
            "rowSpan": 1,
            "text": value,
        }
        for row_index, row in enumerate(normalized_rows)
        for column_index, value in enumerate(row)
        if value
    ]
    if not cells:
        raise DocumentParseError("TABLE_EMPTY")
    return {
        "blockType": "TABLE",
        "cells": cells,
        "columnCount": column_count,
        "locator": locator,
        "ocrConfidence": confidence,
        "readingOrder": 0,
        "rowCount": len(normalized_rows),
    }


def _markdown_table(
    locator: dict[str, object],
    lines: list[str],
    *,
    table_budget: _TableCellBudget,
) -> dict[str, Any]:
    rows: list[list[str]] = []
    for line in lines:
        cells = [cell.strip() for cell in line.strip("|").split("|")]
        if cells and all(re.fullmatch(r":?-{3,}:?", cell) for cell in cells):
            continue
        rows.append(cells)
    return _table(locator, rows, table_budget=table_budget)


def _formula(
    locator: dict[str, object],
    source: str,
    normalized: str,
    confidence: float | None = None,
) -> dict[str, Any]:
    return {
        "blockType": "FORMULA",
        "locator": locator,
        "normalizedFormula": normalized,
        "ocrConfidence": confidence,
        "readingOrder": 0,
        "sourceText": source,
    }


def _normalize_formula(value: str) -> str:
    return re.sub(r"\s+", "", value)


def _ocr_block_to_ir(
    block: OcrBlock,
    locator: dict[str, object],
    *,
    block_budget: _BlockBudget,
    table_budget: _TableCellBudget,
) -> dict[str, Any]:
    if not 0 <= block.confidence <= 1:
        raise DocumentParseError("OCR_CONFIDENCE_INVALID")
    block_budget.reserve()
    if block.block_type == "HEADING" and block.text:
        return _heading(locator, block.text, block.level or 1, block.confidence)
    if block.block_type == "PARAGRAPH" and block.text:
        return _paragraph(locator, block.text, block.confidence)
    if block.block_type == "LIST" and block.items:
        return _list_block(locator, block.items, ordered=block.ordered, confidence=block.confidence)
    if block.block_type == "TABLE" and block.cells and block.row_count and block.column_count:
        # OCR count는 provider-controlled 이므로 dense row allocation 전에 문서 예산을 선점한다.
        table_budget.reserve(
            row_count=block.row_count,
            column_count=block.column_count,
        )
        rows = [["" for _ in range(block.column_count)] for _ in range(block.row_count)]
        for row, column, text in block.cells:
            if row < 0 or row >= block.row_count or column < 0 or column >= block.column_count:
                raise DocumentParseError("OCR_TABLE_INVALID")
            rows[row][column] = text
        return _table(
            locator,
            rows,
            block.confidence,
            table_budget=table_budget,
            pre_reserved=True,
        )
    if block.block_type == "FORMULA" and block.text and block.normalized_formula:
        return _formula(locator, block.text, block.normalized_formula, block.confidence)
    if block.block_type == "CAPTION" and block.text and block.target_reading_order is not None:
        return {
            "blockType": "CAPTION",
            "locator": locator,
            "ocrConfidence": block.confidence,
            "readingOrder": 0,
            "targetReadingOrder": block.target_reading_order,
            "text": block.text,
        }
    raise DocumentParseError("OCR_BLOCK_INVALID")


def _renumber(blocks: list[dict[str, Any]]) -> list[dict[str, Any]]:
    for index, block in enumerate(blocks):
        block["readingOrder"] = index
    return blocks


def _validate_block_bounds(blocks: list[dict[str, Any]], limits: ParserLimits) -> None:
    if not blocks:
        raise DocumentParseError("DOCUMENT_EMPTY")
    if len(blocks) > limits.max_blocks:
        raise DocumentParseError("DOCUMENT_BLOCK_BOUND_EXCEEDED")
    text_characters = 0
    for block in blocks:
        if "text" in block:
            text_characters += len(cast(str, block["text"]))
        if "sourceText" in block:
            text_characters += len(cast(str, block["sourceText"]))
        if "normalizedFormula" in block:
            text_characters += len(cast(str, block["normalizedFormula"]))
        text_characters += sum(len(cast(str, item)) for item in cast(list[object], block.get("items", [])))
        text_characters += sum(
            len(cast(str, cell["text"]))
            for cell in cast(list[dict[str, object]], block.get("cells", []))
        )
    if text_characters <= 0 or text_characters > limits.max_text_characters:
        raise DocumentParseError("DOCUMENT_TEXT_BOUND_EXCEEDED")


def _classify_safety(blocks: list[dict[str, Any]]) -> dict[str, bool]:
    text_parts: list[str] = []
    for block in blocks:
        if isinstance(block.get("text"), str):
            text_parts.append(cast(str, block["text"]))
        if isinstance(block.get("sourceText"), str):
            text_parts.append(cast(str, block["sourceText"]))
        if isinstance(block.get("normalizedFormula"), str):
            text_parts.append(cast(str, block["normalizedFormula"]))
        text_parts.extend(cast(list[str], block.get("items", [])))
        text_parts.extend(
            cast(str, cell["text"])
            for cell in cast(list[dict[str, object]], block.get("cells", []))
        )
    text = "\n".join(text_parts)
    secret = any(pattern.search(text) for pattern in _SECRET_PATTERNS)
    pii = any(pattern.search(text) for pattern in _PII_PATTERNS)
    injection = any(pattern.search(text) for pattern in _PROMPT_INJECTION_PATTERNS)
    return {
        "externalLlmEligible": not (secret or pii or injection),
        "piiDetected": pii,
        "promptInjectionDetected": injection,
        "secretDetected": secret,
    }


def _validate_ocr_identity(backend: OcrBackendPort) -> None:
    if (
        backend.backend not in {"PADDLE_STRUCTURED", "PADDLE_VL", "UNLIMITED_GGUF"}
        or not backend.backend_version
        or len(backend.backend_version) > 128
        or re.fullmatch(r"[0-9a-f]{64}", backend.model_sha256) is None
    ):
        raise DocumentParseError("OCR_BACKEND_IDENTITY_INVALID")


def _extraction_mode(blocks: list[dict[str, Any]], ocr_used: bool) -> str:
    if not ocr_used:
        return "NATIVE"
    native = any(block.get("ocrConfidence") is None for block in blocks)
    return "MIXED" if native else "OCR"
