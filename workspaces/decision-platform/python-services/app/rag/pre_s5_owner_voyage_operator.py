"""Pre-S5 owner 9-format Voyage one-shot의 local-only control plane이다.

author는 문서나 식별자를 manifest에 투영하지 않고 exact plan과 packet digest만 고정한다.
execute는 같은 control을 다시 검증한 뒤 기존 owner Voyage atomic staging 경로를 호출한다.
"""

from __future__ import annotations

import argparse
import hashlib
import json
import os
import re
import secrets
import stat
import sys
from collections.abc import Mapping, Sequence
from dataclasses import dataclass
from datetime import UTC, datetime, timedelta
from pathlib import Path, PurePosixPath

import fitz
import psycopg
from docx import Document as DocxDocument
from openpyxl import Workbook
from PIL import Image, ImageDraw
from pptx import Presentation

from app.rag.bge_acquisition import DEFAULT_MODEL_ROOT
from app.rag.bge_runtime import BgeStaticTokenizer
from app.rag.local_document_parser import (
    LocalDocumentParser,
    OcrBlock,
    OcrPageResult,
)
from app.rag.oa112_downloader import load_oa112_execution_binding
from app.rag.owner_file_io import OwnerFileIoError, read_owner_regular_file
from app.rag.pre_s5_provider_control import (
    PreS5ProviderBinding,
    load_pre_s5_voyage_document_batch_activation,
    resolve_voyage_api_key,
)
from app.rag.pre_s5_voyage_tokenizer import LocalPreS5VoyageContext4Tokenizer
from app.rag.pre_s5_voyage_transport import (
    PreS5VoyageContext4Transport,
    UrllibPreS5VoyageHttpSender,
)
from app.rag.rag_v2_bge_materializer import (
    RagV2OwnerDocumentRequest,
    prepare_owner_document_for_embedding,
)
from app.rag.rag_v2_owner_bge_staging import OwnerBgeStagingMetadata
from app.rag.rag_v2_owner_overlay import PsycopgRagV2OwnerOverlayRepository
from app.rag.rag_v2_owner_voyage_import import (
    OwnerVoyageAttemptLease,
    OwnerVoyageImportPlan,
    PsycopgOwnerVoyageRepository,
    RagV2OwnerVoyageImportExecutor,
    build_owner_voyage_import_item,
    build_owner_voyage_import_plan,
)

OWNER_VOYAGE_SYNTHETIC_FORMATS = (
    "PDF",
    "DOCX",
    "PPTX",
    "XLSX",
    "HTML",
    "MARKDOWN",
    "TXT",
    "PNG",
    "JPEG",
)

_CONTROL_RELATIVE_PATH = "control/owner-voyage-batch-control.v1.json"
_MANIFEST_RELATIVE_PATH = "control/owner-voyage-manifest.v1.json"
_PROFILE_ID = "voyage_context_4_1024_v1"
_CONTROL_FIELDS = {
    "contractId",
    "documents",
    "expiresAt",
    "issuedAt",
    "ownerUserId",
    "schemaVersion",
}
_DOCUMENT_FIELDS = {
    "approvedRoot",
    "documentId",
    "embeddingProfileId",
    "formatId",
    "languageTags",
    "relativePath",
    "retrievalTopics",
    "sanitizedDisplayName",
    "sourceId",
    "sourceRevisionId",
    "ticketId",
}
_EXTENSIONS = (".pdf", ".docx", ".pptx", ".xlsx", ".html", ".md", ".txt", ".png", ".jpg")
_OWNER_ID = re.compile(r"^usr_[a-z0-9][a-z0-9_-]{2,95}$")
_TICKET_ID = re.compile(r"^rti_[0-9a-f]{32}$")
_IDENTIFIER = re.compile(r"^[a-z][a-z0-9_-]{2,127}$")
_SHA256 = re.compile(r"^[0-9a-f]{64}$")
_GIT_OBJECT = re.compile(r"^[0-9a-f]{40,64}$")
_DOCUMENT_PACKET_DIRECTORY = "control/voyage-document-batch-packets"
_PAYMENT_EVIDENCE = "capstone-rag/secrets/voyage-account-evidence/payment-method-on-file.attestation.v1.json"
_OPT_OUT_EVIDENCE = "capstone-rag/secrets/voyage-account-evidence/organization-admin-training-opt-out.attestation.v1.json"


class OwnerVoyageOperatorError(ValueError):
    """owner one-shot control이나 manifest가 closed contract를 위반했다."""


@dataclass(frozen=True, slots=True)
class OwnerVoyageBatchDocument:
    """한 ticket과 한 local source를 결합한 메모리 전용 문서 projection이다."""

    approved_root: Path
    document_id: str
    embedding_profile_id: str
    format_id: str
    language_tags: tuple[str, ...]
    relative_path: str
    retrieval_topics: tuple[str, ...]
    sanitized_display_name: str
    source_id: str
    source_revision_id: str
    import_ticket_id: str


@dataclass(frozen=True, slots=True)
class OwnerVoyageBatchControl:
    """동일 owner의 정확히 9개 synthetic 문서를 one request로 묶는 control이다."""

    owner_user_id: str
    issued_at: datetime
    expires_at: datetime
    documents: tuple[OwnerVoyageBatchDocument, ...]


def load_owner_voyage_batch_control(
    *,
    local_root: Path,
    now: datetime | None = None,
) -> OwnerVoyageBatchControl:
    """고정 0700/0600 control leaf를 읽고 9-format·ticket·profile 불변식을 검증한다."""

    _require_local_root(local_root)
    _require_private_directory(local_root / "control")
    try:
        result = read_owner_regular_file(
            approved_root=local_root,
            relative_path=_CONTROL_RELATIVE_PATH,
            max_bytes=64 * 1024,
        )
    except OwnerFileIoError as error:
        raise OwnerVoyageOperatorError("OWNER_VOYAGE_BATCH_CONTROL_BOUNDARY") from error
    try:
        value = json.loads(
            result.content.decode("utf-8", errors="strict"),
            object_pairs_hook=_reject_duplicate_keys,
        )
    except (UnicodeDecodeError, ValueError, json.JSONDecodeError) as error:
        raise OwnerVoyageOperatorError("OWNER_VOYAGE_BATCH_CONTROL_INVALID") from error
    if not isinstance(value, dict) or set(value) != _CONTROL_FIELDS:
        raise OwnerVoyageOperatorError("OWNER_VOYAGE_BATCH_CONTROL_INVALID")
    try:
        issued_at = _parse_instant(value["issuedAt"])
        expires_at = _parse_instant(value["expiresAt"])
    except (KeyError, TypeError, ValueError) as error:
        raise OwnerVoyageOperatorError("OWNER_VOYAGE_BATCH_CONTROL_INVALID") from error
    current = (now or datetime.now(UTC)).astimezone(UTC)
    owner_user_id = value.get("ownerUserId")
    raw_documents = value.get("documents")
    if (
        value.get("contractId") != "pre-s5-owner-voyage-batch-control/v1"
        or value.get("schemaVersion") != 1
        or not isinstance(owner_user_id, str)
        or _OWNER_ID.fullmatch(owner_user_id) is None
        or not isinstance(raw_documents, list)
        or len(raw_documents) != len(OWNER_VOYAGE_SYNTHETIC_FORMATS)
        or expires_at <= issued_at
        or expires_at - issued_at > timedelta(minutes=5)
        or current < issued_at
        or current >= expires_at
    ):
        raise OwnerVoyageOperatorError("OWNER_VOYAGE_BATCH_CONTROL_INVALID")

    documents = tuple(
        _parse_document(raw, expected_format=format_id, expected_extension=extension)
        for raw, format_id, extension in zip(
            raw_documents,
            OWNER_VOYAGE_SYNTHETIC_FORMATS,
            _EXTENSIONS,
            strict=True,
        )
    )
    ticket_ids = {document.import_ticket_id for document in documents}
    source_ids = {document.source_id for document in documents}
    revision_ids = {document.source_revision_id for document in documents}
    document_ids = {document.document_id for document in documents}
    if any(len(values) != len(documents) for values in (ticket_ids, source_ids, revision_ids, document_ids)):
        raise OwnerVoyageOperatorError("OWNER_VOYAGE_BATCH_CONTROL_INVALID")
    return OwnerVoyageBatchControl(
        owner_user_id=owner_user_id,
        issued_at=issued_at,
        expires_at=expires_at,
        documents=documents,
    )


def author_owner_voyage_manifest(
    *,
    output_path: Path,
    head_commit: str,
    tree_object: str,
    ci_digest: str,
    security_digest: str,
    plan_sha256: str,
    batch_manifest_sha256: str,
    owner_scope_sha256: str,
    ticket_set_sha256: str,
    tokenizer_sha256: str,
    packet_sha256: str,
    document_count: int,
    chunk_count: int,
    token_count: int,
    expires_at: datetime,
) -> str:
    """raw owner identity 없이 exact one-call approval manifest를 0600으로 기록한다."""

    if (
        not isinstance(output_path, Path)
        or not output_path.is_absolute()
        or _GIT_OBJECT.fullmatch(head_commit) is None
        or _GIT_OBJECT.fullmatch(tree_object) is None
        or any(
            _SHA256.fullmatch(value) is None
            for value in (
                ci_digest,
                security_digest,
                plan_sha256,
                batch_manifest_sha256,
                owner_scope_sha256,
                ticket_set_sha256,
                tokenizer_sha256,
                packet_sha256,
            )
        )
        or document_count != 9
        or not 1 <= chunk_count <= 16_000
        or not 1 <= token_count <= 55_000
        or not isinstance(expires_at, datetime)
        or expires_at.tzinfo is None
        or not datetime.now(UTC) < expires_at.astimezone(UTC) <= datetime.now(UTC) + timedelta(hours=2)
    ):
        raise OwnerVoyageOperatorError("OWNER_VOYAGE_MANIFEST_INVALID")
    payload: dict[str, object] = {
        "approvalScope": "PRE_S5_OWNER_VOYAGE_SYNTHETIC_ONE_SHOT",
        "batchManifestSha256": batch_manifest_sha256,
        "binding": {
            "ciDigest": ci_digest,
            "headCommit": head_commit,
            "securityDigest": security_digest,
            "treeObject": tree_object,
        },
        "chunkCount": chunk_count,
        "documentCount": document_count,
        "embeddingProfileId": _PROFILE_ID,
        "expiresAt": _format_instant(expires_at),
        "operation": "OWNER_VOYAGE_DOCUMENT_IMPORT",
        "ownerScopeSha256": owner_scope_sha256,
        "packetSha256": packet_sha256,
        "physicalCallCap": 1,
        "planSha256": plan_sha256,
        "rawArtifactCount": 0,
        "retryCount": 0,
        "schemaVersion": 1,
        "ticketSetSha256": ticket_set_sha256,
        "tokenCount": token_count,
        "tokenizerSha256": tokenizer_sha256,
    }
    content = _canonical_json(payload)
    _write_private_json(output_path, content)
    return hashlib.sha256(content).hexdigest()


def main(argv: Sequence[str] | None = None) -> int:
    """`author|execute`만 노출하고 raw path·secret을 argv나 stdout에 싣지 않는다."""

    parser = argparse.ArgumentParser(prog="pre-s5-owner-voyage")
    parser.add_argument("command", choices=("author", "execute"))
    arguments = parser.parse_args(tuple(sys.argv[1:] if argv is None else argv))
    try:
        local_root = _environment_absolute_path("CAPSTONE_RAG_LOCAL_ROOT")
        _require_local_root(local_root)
        if arguments.command == "author":
            _prepare_synthetic_batch(local_root=local_root)
        control = load_owner_voyage_batch_control(local_root=local_root)
        plan = _materialize_plan(control=control, local_root=local_root)
        binding = _execution_binding(local_root=local_root)
        if arguments.command == "author":
            packet_sha256 = _author_document_packet(
                local_root=local_root,
                control=control,
                plan=plan,
                binding=binding,
            )
            manifest_sha256 = author_owner_voyage_manifest(
                output_path=local_root / _MANIFEST_RELATIVE_PATH,
                head_commit=binding.head_commit,
                tree_object=binding.tree_object,
                ci_digest=binding.ci_digest,
                security_digest=binding.security_digest,
                plan_sha256=plan.plan_sha256,
                batch_manifest_sha256=plan.batch.batch_manifest_sha256,
                owner_scope_sha256=plan.owner_scope_sha256,
                ticket_set_sha256=plan.ticket_set_sha256,
                tokenizer_sha256=plan.tokenizer_sha256,
                packet_sha256=packet_sha256,
                document_count=len(plan.items),
                chunk_count=plan.batch.chunk_count,
                token_count=plan.batch.token_count,
                expires_at=control.expires_at,
            )
            _emit(
                {
                    "chunkCount": plan.batch.chunk_count,
                    "documentCount": len(plan.items),
                    "manifestSha256": manifest_sha256,
                    "physicalCallCap": 1,
                    "planSha256": plan.plan_sha256,
                    "providerCalls": 0,
                    "state": "AUTHORED",
                    "tokenCount": plan.batch.token_count,
                }
            )
            return 0
        return _execute(local_root=local_root, control=control, plan=plan, binding=binding)
    except Exception as error:
        code = str(error)
        if not code or len(code) > 160 or not re.fullmatch(r"[A-Z0-9_]+", code):
            code = "OWNER_VOYAGE_OPERATOR_FAILED"
        return _emit_failure(code)


def _materialize_plan(
    *,
    control: OwnerVoyageBatchControl,
    local_root: Path,
) -> OwnerVoyageImportPlan:
    """각 control 문서를 정확히 한 번 parse해 one-request plan으로 결합한다."""

    tokenizer_sha256 = _required_sha256_environment("CAPSTONE_RAG_VOYAGE_TOKENIZER_SHA256")
    voyage_root = _optional_environment_absolute_path("CAPSTONE_RAG_VOYAGE_TOKENIZER_ROOT") or local_root
    voyage_tokenizer = LocalPreS5VoyageContext4Tokenizer.from_local_root(
        local_root=voyage_root,
        expected_sha256=tokenizer_sha256,
    )
    bge_root = _optional_environment_absolute_path("CAPSTONE_RAG_BGE_PACKET_ROOT") or DEFAULT_MODEL_ROOT
    canonical_tokenizer = BgeStaticTokenizer.from_file(bge_root / "onnx" / "tokenizer.json")
    parser = LocalDocumentParser(ocr_backend=_SyntheticOwnerOcrBackend())
    items = []
    for document in control.documents:
        prepared = prepare_owner_document_for_embedding(
            parser=parser,
            tokenizer=canonical_tokenizer,
            request=RagV2OwnerDocumentRequest(
                approved_root=document.approved_root,
                relative_path=document.relative_path,
                document_id=document.document_id,
                source_id=document.source_id,
                source_revision_id=document.source_revision_id,
                language_tags=document.language_tags,
                embedding_profile_id=document.embedding_profile_id,
            ),
            # reserve 함수가 outbound 직전에 current v2 consent와 ticket을 다시 확인한다.
            external_processing_authorized=True,
        )
        items.append(
            build_owner_voyage_import_item(
                import_ticket_id=document.import_ticket_id,
                prepared=prepared,
                metadata=OwnerBgeStagingMetadata(
                    sanitized_display_name=document.sanitized_display_name,
                    retrieval_topics=document.retrieval_topics,
                ),
                tokenizer_version=f"voyage-context-4-{tokenizer_sha256[:16]}",
            )
        )
    return build_owner_voyage_import_plan(
        owner_user_id=control.owner_user_id,
        items=tuple(items),
        token_counter=voyage_tokenizer,
    )


def _execution_binding(*, local_root: Path) -> PreS5ProviderBinding:
    relative_path = os.environ.get(
        "CAPSTONE_RAG_OWNER_VOYAGE_EXECUTION_EVIDENCE",
        "pre-s5-voyage-execution-evidence.v1.json",
    ).strip()
    evidence = load_oa112_execution_binding(
        approved_root=local_root,
        relative_path=relative_path,
        repository_root=Path(__file__).resolve().parents[5],
    )
    return PreS5ProviderBinding(
        head_commit=evidence.head_sha,
        tree_object=evidence.tree_sha256,
        ci_digest=evidence.ci_digest,
        security_digest=evidence.security_digest,
    )


def _author_document_packet(
    *,
    local_root: Path,
    control: OwnerVoyageBatchControl,
    plan: OwnerVoyageImportPlan,
    binding: PreS5ProviderBinding,
) -> str:
    """기존 strict provider loader와 byte-for-byte 호환되는 one-shot child packet을 만든다."""

    repository_root = Path(__file__).resolve().parents[5]
    payment_hash = _hash_private_evidence(repository_root / _PAYMENT_EVIDENCE)
    opt_out_hash = _hash_private_evidence(repository_root / _OPT_OUT_EVIDENCE)
    rate_hash = _required_sha256_environment("PRE_S5_VOYAGE_RATE_EVIDENCE_SHA256")
    operator = os.environ.get("PRE_S5_OPERATOR", "").strip()
    if re.fullmatch(r"[a-z0-9][a-z0-9._@-]{2,127}", operator) is None:
        raise OwnerVoyageOperatorError("OWNER_VOYAGE_OPERATOR_INVALID")
    try:
        price = int(os.environ.get("PRE_S5_VOYAGE_INPUT_MICROUSD_PER_TOKEN", ""))
    except ValueError as error:
        raise OwnerVoyageOperatorError("OWNER_VOYAGE_PRICING_INVALID") from error
    if not 1 <= price <= 1_000_000 or plan.batch.estimated_response_bytes > 16 * 1024 * 1024:
        raise OwnerVoyageOperatorError("OWNER_VOYAGE_PRICING_INVALID")
    issued_at = datetime.now(UTC)
    expires_at = min(control.expires_at, issued_at + timedelta(hours=2))
    if expires_at <= issued_at:
        raise OwnerVoyageOperatorError("OWNER_VOYAGE_BATCH_CONTROL_INVALID")
    packet: dict[str, object] = {
        "batchCount": plan.batch.batch_count,
        "batchId": plan.batch.batch_id,
        "batchManifestSha256": plan.batch.batch_manifest_sha256,
        "batchOrdinal": plan.batch.batch_ordinal,
        "batchPlanSha256": plan.plan_sha256,
        "byteCap": plan.batch.estimated_response_bytes,
        "chunkCount": plan.batch.chunk_count,
        "ciDigest": binding.ci_digest,
        "costCapMicrousd": 55_000 * price,
        "date": "NONE",
        "endpoint": "/v1/contextualizedembeddings",
        "expiresAt": _format_instant(expires_at),
        "groupCount": plan.batch.group_count,
        "headCommit": binding.head_commit,
        "inputMicrousdPerToken": price,
        "issuedAt": _format_instant(issued_at),
        "logicalCallCap": 1,
        "nonce": f"ps5_owner_{secrets.token_hex(16)}",
        "operation": "CONTEXTUALIZED_DOCUMENT_EMBEDDING",
        "operator": operator,
        "organizationTrainingOptOutEvidenceSha256": opt_out_hash,
        "origin": "https://api.voyageai.com",
        "paymentMethodPrivacyEvidenceSha256": payment_hash,
        "physicalCallCap": 1,
        "provider": "VOYAGE",
        "query": "MANIFEST_BOUND_ORDERED_PRECHUNKED_DOCUMENT_BATCH",
        "rateEvidenceSha256": rate_hash,
        "rawArtifactCount": 0,
        "retryCount": 0,
        "schemaVersion": "pre-s5-voyage-document-batch-activation/v1",
        "securityDigest": binding.security_digest,
        "state": "APPROVED",
        "symbol": "NONE",
        "tokenCap": 55_000,
        "tokenCount": plan.batch.token_count,
        "tokenizerSha256": plan.tokenizer_sha256,
        "treeObject": binding.tree_object,
    }
    content = _canonical_json(packet)
    packet_path = local_root / _DOCUMENT_PACKET_DIRECTORY / f"{plan.batch.batch_id}.json"
    _write_private_json(packet_path, content)
    # production loader로 즉시 self-check해 author와 execute validator drift를 제거한다.
    activation = load_pre_s5_voyage_document_batch_activation(
        local_root=local_root,
        binding=binding,
        batch_plan_sha256=plan.plan_sha256,
        batch_id=plan.batch.batch_id,
        batch_manifest_sha256=plan.batch.batch_manifest_sha256,
        batch_ordinal=plan.batch.batch_ordinal,
        batch_count=plan.batch.batch_count,
        token_count=plan.batch.token_count,
        chunk_count=plan.batch.chunk_count,
        group_count=plan.batch.group_count,
        estimated_response_bytes=plan.batch.estimated_response_bytes,
    )
    return activation.packet_sha256


def _execute(
    *,
    local_root: Path,
    control: OwnerVoyageBatchControl,
    plan: OwnerVoyageImportPlan,
    binding: PreS5ProviderBinding,
) -> int:
    """승인 manifest와 child packet이 일치할 때만 정확히 한 provider call을 수행한다."""

    from app.rag.content_cli import _load_owner_voyage_manifest

    manifest_sha256, manifest = _load_owner_voyage_manifest(
        local_root=local_root,
        plan=plan,
        binding=binding,
    )
    activation = load_pre_s5_voyage_document_batch_activation(
        local_root=local_root,
        binding=binding,
        batch_plan_sha256=plan.plan_sha256,
        batch_id=plan.batch.batch_id,
        batch_manifest_sha256=plan.batch.batch_manifest_sha256,
        batch_ordinal=plan.batch.batch_ordinal,
        batch_count=plan.batch.batch_count,
        token_count=plan.batch.token_count,
        chunk_count=plan.batch.chunk_count,
        group_count=plan.batch.group_count,
        estimated_response_bytes=plan.batch.estimated_response_bytes,
    )
    if manifest.get("packetSha256") != activation.packet_sha256:
        raise OwnerVoyageOperatorError("OWNER_VOYAGE_MANIFEST_BINDING")
    writer_dsn = os.environ.get("CAPSTONE_RAG_WRITER_DATABASE_DSN", "").strip()
    admin_dsn = os.environ.get("CAPSTONE_RAG_ADMIN_DATABASE_DSN", "").strip()
    if not writer_dsn or not admin_dsn:
        raise OwnerVoyageOperatorError("OWNER_VOYAGE_DATABASE_DSN")
    repository = PsycopgOwnerVoyageRepository(database_dsn=writer_dsn)
    lease = OwnerVoyageAttemptLease(
        repository=repository,
        plan=plan,
        packet_sha256=activation.packet_sha256,
        approval_manifest_sha256=manifest_sha256,
        nonce_sha256=activation.nonce_sha256,
    )
    voyage_root = _optional_environment_absolute_path("CAPSTONE_RAG_VOYAGE_TOKENIZER_ROOT") or local_root
    voyage_tokenizer = LocalPreS5VoyageContext4Tokenizer.from_local_root(
        local_root=voyage_root,
        expected_sha256=plan.tokenizer_sha256,
    )
    transport = PreS5VoyageContext4Transport(
        activation=activation,
        api_key=resolve_voyage_api_key(os.environ),
        lease=lease,
        token_counter=voyage_tokenizer,
        sender=UrllibPreS5VoyageHttpSender(),
    )
    receipt = RagV2OwnerVoyageImportExecutor(repository=lease).execute(
        plan=plan,
        transport=transport,
    )
    overlay = PsycopgRagV2OwnerOverlayRepository(database_dsn=admin_dsn).prepare_and_activate(
        owner_user_id=control.owner_user_id
    )
    _emit(
        {
            "chunkCount": receipt.chunk_count,
            "documentCount": receipt.document_count,
            "embeddingProfileId": receipt.embedding_profile_id,
            "ownerBundleState": overlay.state,
            "physicalCalls": 1,
            "state": receipt.state,
        }
    )
    return 0


class _SyntheticOwnerOcrBackend:
    """synthetic PNG/JPEG 두 fixture만 network 없이 deterministic text로 투영한다."""

    backend = "PADDLE_STRUCTURED"
    backend_version = "pre-s5-synthetic-fixture-v1"
    model_sha256 = hashlib.sha256(b"pre-s5-synthetic-owner-ocr-fixture-v1").hexdigest()

    def parse_page(self, *, png_bytes: bytes, page_number: int) -> OcrPageResult:
        if not png_bytes.startswith(b"\x89PNG\r\n\x1a\n") or page_number != 1:
            raise OwnerVoyageOperatorError("OWNER_VOYAGE_SYNTHETIC_OCR_INVALID")
        return OcrPageResult(
            blocks=(
                OcrBlock(
                    block_type="PARAGRAPH",
                    confidence=1.0,
                    text="분산투자는 서로 다른 자산의 위험 노출을 나누어 집중 위험을 낮춘다.",
                ),
            )
        )


def _prepare_synthetic_batch(*, local_root: Path) -> None:
    """safe 9-format fixture와 fresh v2 consent/tickets를 local DB authority로 즉시 만든다."""

    owner_user_id = os.environ.get("PRE_S5_OWNER_SYNTHETIC_USER_ID", "").strip()
    app_dsn = os.environ.get("CAPSTONE_RAG_APP_DATABASE_DSN", "").strip()
    if _OWNER_ID.fullmatch(owner_user_id) is None or not app_dsn:
        raise OwnerVoyageOperatorError("OWNER_VOYAGE_SYNTHETIC_CONFIGURATION")
    nonce = secrets.token_hex(16)
    synthetic_root = local_root / "owner-synthetic"
    _create_or_require_private_directory(synthetic_root)
    document_root = synthetic_root / nonce
    document_root.mkdir(mode=0o700, parents=False, exist_ok=False)
    os.chmod(document_root, 0o700)
    base_text = (
        "Diversification combines assets with different risk exposures. "
        "Lower correlation can reduce portfolio volatility while expected returns remain uncertain."
    )
    filenames = (
        "safe-01.pdf",
        "safe-02.docx",
        "safe-03.pptx",
        "safe-04.xlsx",
        "safe-05.html",
        "safe-06.md",
        "safe-07.txt",
        "safe-08.png",
        "safe-09.jpg",
    )
    _write_synthetic_documents(document_root, filenames=filenames, text=base_text)
    issued_at = datetime.now(UTC)
    documents: list[dict[str, object]] = []
    ticket_ids = tuple(f"rti_{secrets.token_hex(16)}" for _ in filenames)
    expires_at: datetime | None = None
    internal_consent_id = f"cns_v2_{secrets.token_hex(16)}"
    public_consent_id = f"rce_pres5_{secrets.token_hex(16)}"
    disclosure_digest = hashlib.sha256(b"pre-s5-owner-voyage-synthetic-disclosure-v1").hexdigest()
    policy_digest = hashlib.sha256(b"EXTERNAL_AI_RAG_V2").hexdigest()
    processor_digest = hashlib.sha256(b"VOYAGE_CONTEXT_4_NO_TRAINING").hexdigest()
    try:
        with psycopg.connect(app_dsn, autocommit=False, connect_timeout=2) as connection:
            with connection.transaction():
                connection.execute("SELECT set_config('app.actor_user_id', %s, true)", (owner_user_id,))
                connection.execute(
                    """
                    SELECT public.record_rag_v2_immutable_consent_v2(
                      %s, %s, %s, 'GRANT', %s, %s, %s
                    )
                    """,
                    (
                        owner_user_id,
                        internal_consent_id,
                        public_consent_id,
                        disclosure_digest,
                        policy_digest,
                        processor_digest,
                    ),
                ).fetchone()
                for ticket_id in ticket_ids:
                    row = connection.execute(
                        """
                        SELECT public.issue_rag_v2_immutable_import_ticket_v2(
                          %s, %s, 'OWNER_IMPORT', 'RAG_V2_OWNER_DOCUMENT_V2', %s
                        )
                        """,
                        (owner_user_id, ticket_id, _PROFILE_ID),
                    ).fetchone()
                    if row is None or len(row) != 1 or not isinstance(row[0], datetime):
                        raise OwnerVoyageOperatorError("OWNER_VOYAGE_SYNTHETIC_TICKET")
                    candidate = row[0].astimezone(UTC)
                    expires_at = candidate if expires_at is None else min(expires_at, candidate)
    except OwnerVoyageOperatorError:
        raise
    except psycopg.Error as error:
        raise OwnerVoyageOperatorError("OWNER_VOYAGE_SYNTHETIC_TICKET") from error
    if expires_at is None or expires_at <= issued_at:
        raise OwnerVoyageOperatorError("OWNER_VOYAGE_SYNTHETIC_TICKET")
    for index, (format_id, filename, ticket_id) in enumerate(
        zip(OWNER_VOYAGE_SYNTHETIC_FORMATS, filenames, ticket_ids, strict=True),
        start=1,
    ):
        documents.append(
            {
                "approvedRoot": str(document_root),
                "documentId": f"doc_owner_voyage_synthetic_{nonce}_{index:02d}",
                "embeddingProfileId": _PROFILE_ID,
                "formatId": format_id,
                "languageTags": ["en"],
                "relativePath": filename,
                "retrievalTopics": ["RISK", "FINANCIAL_ENGINEERING"],
                "sanitizedDisplayName": f"Synthetic diversification evidence {index:02d}",
                "sourceId": f"src_owner_voyage_synthetic_{nonce}_{index:02d}",
                "sourceRevisionId": f"srv_owner_voyage_synthetic_{nonce}_{index:02d}",
                "ticketId": ticket_id,
            }
        )
    control: dict[str, object] = {
        "contractId": "pre-s5-owner-voyage-batch-control/v1",
        "documents": documents,
        "expiresAt": _format_instant(expires_at),
        "issuedAt": _format_instant(issued_at),
        "ownerUserId": owner_user_id,
        "schemaVersion": 1,
    }
    _write_private_json(local_root / _CONTROL_RELATIVE_PATH, _canonical_json(control))


def _write_synthetic_documents(
    root: Path,
    *,
    filenames: tuple[str, ...],
    text: str,
) -> None:
    """format별 native writer를 써서 실제 parser 경로가 읽을 최소 안전 fixture를 만든다."""

    pdf = fitz.open()
    page = pdf.new_page()
    page.insert_text((72, 72), text)
    _write_private_bytes(root / filenames[0], pdf.tobytes())
    pdf.close()

    document = DocxDocument()
    document.add_heading("Diversification", level=1)
    document.add_paragraph(text)
    _save_private_via_temporary(document.save, root / filenames[1])

    presentation = Presentation()
    slide = presentation.slides.add_slide(presentation.slide_layouts[1])
    slide.shapes.title.text = "Diversification"
    slide.placeholders[1].text = text
    _save_private_via_temporary(presentation.save, root / filenames[2])

    workbook = Workbook()
    sheet = workbook.active
    if sheet is None:  # pragma: no cover - openpyxl의 새 workbook은 항상 active sheet를 가진다.
        raise OwnerVoyageOperatorError("OWNER_VOYAGE_SYNTHETIC_FILE")
    sheet.title = "Evidence"
    sheet.append(("Topic", "Evidence"))
    sheet.append(("Diversification", text))
    _save_private_via_temporary(workbook.save, root / filenames[3])
    workbook.close()

    _write_private_bytes(
        root / filenames[4],
        f"<!doctype html><html><body><h1>Diversification</h1><p>{text}</p></body></html>".encode(),
    )
    _write_private_bytes(root / filenames[5], f"# Diversification\n\n{text}\n".encode())
    _write_private_bytes(root / filenames[6], f"Diversification\n{text}\n".encode())

    for filename, image_format in ((filenames[7], "PNG"), (filenames[8], "JPEG")):
        image = Image.new("RGB", (1200, 240), "white")
        draw = ImageDraw.Draw(image)
        draw.text((30, 80), "Diversification reduces concentration risk.", fill="black")
        temporary = root / f".{filename}.image.tmp"
        image.save(temporary, format=image_format)
        image.close()
        _publish_private_temporary(temporary, root / filename)


def _save_private_via_temporary(writer: object, target: Path) -> None:
    temporary = target.with_name(f".{target.name}.writer.tmp")
    try:
        writer(str(temporary))  # type: ignore[operator]
        _publish_private_temporary(temporary, target)
    except OwnerVoyageOperatorError:
        raise
    except Exception as error:
        raise OwnerVoyageOperatorError("OWNER_VOYAGE_SYNTHETIC_FILE") from error


def _publish_private_temporary(temporary: Path, target: Path) -> None:
    try:
        metadata = temporary.lstat()
        if not stat.S_ISREG(metadata.st_mode) or metadata.st_nlink != 1 or metadata.st_size <= 0:
            raise OwnerVoyageOperatorError("OWNER_VOYAGE_SYNTHETIC_FILE")
        os.chmod(temporary, 0o600)
        os.replace(temporary, target)
    except OwnerVoyageOperatorError:
        raise
    except OSError as error:
        raise OwnerVoyageOperatorError("OWNER_VOYAGE_SYNTHETIC_FILE") from error


def _write_private_bytes(path: Path, content: bytes) -> None:
    if not content or len(content) > 100 * 1024 * 1024:
        raise OwnerVoyageOperatorError("OWNER_VOYAGE_SYNTHETIC_FILE")
    descriptor = -1
    try:
        descriptor = os.open(
            path,
            os.O_WRONLY | os.O_CREAT | os.O_EXCL | os.O_NOFOLLOW | os.O_CLOEXEC,
            0o600,
        )
        os.fchmod(descriptor, 0o600)
        offset = 0
        while offset < len(content):
            written = os.write(descriptor, content[offset:])
            if written <= 0:
                raise OwnerVoyageOperatorError("OWNER_VOYAGE_SYNTHETIC_FILE")
            offset += written
        os.fsync(descriptor)
    except OwnerVoyageOperatorError:
        raise
    except OSError as error:
        raise OwnerVoyageOperatorError("OWNER_VOYAGE_SYNTHETIC_FILE") from error
    finally:
        if descriptor >= 0:
            os.close(descriptor)


def _parse_document(
    value: object,
    *,
    expected_format: str,
    expected_extension: str,
) -> OwnerVoyageBatchDocument:
    if not isinstance(value, dict) or set(value) != _DOCUMENT_FIELDS:
        raise OwnerVoyageOperatorError("OWNER_VOYAGE_BATCH_CONTROL_INVALID")
    approved_root = value.get("approvedRoot")
    relative_path = value.get("relativePath")
    language_tags = value.get("languageTags")
    retrieval_topics = value.get("retrievalTopics")
    scalar_names = (
        "documentId",
        "sanitizedDisplayName",
        "sourceId",
        "sourceRevisionId",
        "ticketId",
    )
    if (
        not isinstance(approved_root, str)
        or not Path(approved_root).is_absolute()
        or ".." in Path(approved_root).parts
        or not isinstance(relative_path, str)
        or not _safe_relative_path(relative_path)
        or PurePosixPath(relative_path).suffix.lower() != expected_extension
        or value.get("formatId") != expected_format
        or value.get("embeddingProfileId") != _PROFILE_ID
        or not isinstance(language_tags, list)
        or not language_tags
        or any(not isinstance(item, str) or not 1 <= len(item) <= 35 for item in language_tags)
        or not isinstance(retrieval_topics, list)
        or not retrieval_topics
        or any(not isinstance(item, str) or not 1 <= len(item) <= 64 for item in retrieval_topics)
        or any(not isinstance(value.get(name), str) for name in scalar_names)
        or _TICKET_ID.fullmatch(str(value.get("ticketId"))) is None
        or any(
            _IDENTIFIER.fullmatch(str(value.get(name))) is None
            for name in ("documentId", "sourceId", "sourceRevisionId")
        )
        or not 1 <= len(str(value.get("sanitizedDisplayName"))) <= 160
    ):
        raise OwnerVoyageOperatorError("OWNER_VOYAGE_BATCH_CONTROL_INVALID")
    return OwnerVoyageBatchDocument(
        approved_root=Path(approved_root),
        document_id=str(value["documentId"]),
        embedding_profile_id=_PROFILE_ID,
        format_id=expected_format,
        language_tags=tuple(language_tags),
        relative_path=relative_path,
        retrieval_topics=tuple(retrieval_topics),
        sanitized_display_name=str(value["sanitizedDisplayName"]),
        source_id=str(value["sourceId"]),
        source_revision_id=str(value["sourceRevisionId"]),
        import_ticket_id=str(value["ticketId"]),
    )


def _require_local_root(path: Path) -> None:
    if not isinstance(path, Path) or not path.is_absolute() or ".." in path.parts:
        raise OwnerVoyageOperatorError("OWNER_VOYAGE_BATCH_CONTROL_BOUNDARY")
    _require_private_directory(path)


def _require_private_directory(path: Path) -> None:
    try:
        metadata = path.lstat()
    except OSError as error:
        raise OwnerVoyageOperatorError("OWNER_VOYAGE_BATCH_CONTROL_BOUNDARY") from error
    if (
        not stat.S_ISDIR(metadata.st_mode)
        or metadata.st_uid != os.geteuid()
        or stat.S_IMODE(metadata.st_mode) != 0o700
    ):
        raise OwnerVoyageOperatorError("OWNER_VOYAGE_BATCH_CONTROL_BOUNDARY")


def _write_private_json(path: Path, content: bytes) -> None:
    if not path.is_absolute() or ".." in path.parts or not content or len(content) > 64 * 1024:
        raise OwnerVoyageOperatorError("OWNER_VOYAGE_MANIFEST_BOUNDARY")
    _create_or_require_private_directory(path.parent, code="OWNER_VOYAGE_MANIFEST_BOUNDARY")
    temporary_name = f".{path.name}.{os.getpid()}.tmp"
    parent_descriptor = -1
    descriptor = -1
    try:
        parent_descriptor = os.open(
            path.parent,
            os.O_RDONLY | os.O_DIRECTORY | os.O_NOFOLLOW | os.O_CLOEXEC,
        )
        _validate_private_directory_metadata(
            os.fstat(parent_descriptor),
            code="OWNER_VOYAGE_MANIFEST_BOUNDARY",
        )
        descriptor = os.open(
            temporary_name,
            os.O_WRONLY | os.O_CREAT | os.O_EXCL | os.O_NOFOLLOW | os.O_CLOEXEC,
            0o600,
            dir_fd=parent_descriptor,
        )
        os.fchmod(descriptor, 0o600)
        offset = 0
        while offset < len(content):
            written = os.write(descriptor, content[offset:])
            if written <= 0:
                raise OwnerVoyageOperatorError("OWNER_VOYAGE_MANIFEST_BOUNDARY")
            offset += written
        os.fsync(descriptor)
        metadata = os.fstat(descriptor)
        if not stat.S_ISREG(metadata.st_mode) or metadata.st_nlink != 1 or metadata.st_size != len(content):
            raise OwnerVoyageOperatorError("OWNER_VOYAGE_MANIFEST_BOUNDARY")
        os.close(descriptor)
        descriptor = -1
        try:
            existing = os.stat(path.name, dir_fd=parent_descriptor, follow_symlinks=False)
        except FileNotFoundError:
            existing = None
        if existing is not None:
            if (
                not stat.S_ISREG(existing.st_mode)
                or existing.st_uid != os.geteuid()
                or existing.st_nlink != 1
                or stat.S_IMODE(existing.st_mode) != 0o600
            ):
                raise OwnerVoyageOperatorError("OWNER_VOYAGE_MANIFEST_BOUNDARY")
        os.replace(
            temporary_name,
            path.name,
            src_dir_fd=parent_descriptor,
            dst_dir_fd=parent_descriptor,
        )
        os.chmod(path.name, 0o600, dir_fd=parent_descriptor, follow_symlinks=False)
    except OwnerVoyageOperatorError:
        raise
    except OSError as error:
        raise OwnerVoyageOperatorError("OWNER_VOYAGE_MANIFEST_BOUNDARY") from error
    finally:
        if descriptor >= 0:
            os.close(descriptor)
        if parent_descriptor >= 0:
            try:
                os.unlink(temporary_name, dir_fd=parent_descriptor)
            except FileNotFoundError:
                pass
            os.close(parent_descriptor)


def _create_or_require_private_directory(
    path: Path,
    *,
    code: str = "OWNER_VOYAGE_BATCH_CONTROL_BOUNDARY",
) -> None:
    """trusted 0700 parent 아래 한 단계만 만들고 link traversal을 막는다."""

    try:
        metadata = path.lstat()
    except FileNotFoundError:
        _require_private_directory_with_code(path.parent, code=code)
        try:
            path.mkdir(mode=0o700, parents=False, exist_ok=False)
        except OSError as error:
            raise OwnerVoyageOperatorError(code) from error
        metadata = path.lstat()
    except OSError as error:
        raise OwnerVoyageOperatorError(code) from error
    _validate_private_directory_metadata(metadata, code=code)


def _require_private_directory_with_code(path: Path, *, code: str) -> None:
    try:
        metadata = path.lstat()
    except OSError as error:
        raise OwnerVoyageOperatorError(code) from error
    _validate_private_directory_metadata(metadata, code=code)


def _validate_private_directory_metadata(metadata: os.stat_result, *, code: str) -> None:
    if (
        not stat.S_ISDIR(metadata.st_mode)
        or metadata.st_uid != os.geteuid()
        or stat.S_IMODE(metadata.st_mode) != 0o700
    ):
        raise OwnerVoyageOperatorError(code)


def _safe_relative_path(value: str) -> bool:
    if not value or value.startswith("/") or value.endswith("/") or "\\" in value or "\x00" in value:
        return False
    return all(part not in {"", ".", ".."} and ":" not in part for part in PurePosixPath(value).parts)


def _environment_absolute_path(name: str) -> Path:
    value = os.environ.get(name, "").strip()
    path = Path(value)
    if not value or not path.is_absolute() or ".." in path.parts:
        raise OwnerVoyageOperatorError("OWNER_VOYAGE_ENVIRONMENT")
    return path


def _optional_environment_absolute_path(name: str) -> Path | None:
    value = os.environ.get(name, "").strip()
    if not value:
        return None
    path = Path(value)
    if not path.is_absolute() or ".." in path.parts:
        raise OwnerVoyageOperatorError("OWNER_VOYAGE_ENVIRONMENT")
    return path


def _required_sha256_environment(name: str) -> str:
    value = os.environ.get(name, "").strip()
    if _SHA256.fullmatch(value) is None:
        raise OwnerVoyageOperatorError("OWNER_VOYAGE_ENVIRONMENT")
    return value


def _hash_private_evidence(path: Path) -> str:
    descriptor = -1
    try:
        metadata = path.lstat()
        if (
            not stat.S_ISREG(metadata.st_mode)
            or metadata.st_uid != os.geteuid()
            or metadata.st_nlink != 1
            or stat.S_IMODE(metadata.st_mode) != 0o600
            or not 1 <= metadata.st_size <= 64 * 1024
        ):
            raise OwnerVoyageOperatorError("OWNER_VOYAGE_EVIDENCE_BOUNDARY")
        descriptor = os.open(path, os.O_RDONLY | os.O_NOFOLLOW | os.O_CLOEXEC)
        opened = os.fstat(descriptor)
        if (
            (opened.st_dev, opened.st_ino) != (metadata.st_dev, metadata.st_ino)
            or not stat.S_ISREG(opened.st_mode)
            or opened.st_uid != os.geteuid()
            or opened.st_nlink != 1
            or stat.S_IMODE(opened.st_mode) != 0o600
            or opened.st_size != metadata.st_size
        ):
            raise OwnerVoyageOperatorError("OWNER_VOYAGE_EVIDENCE_BOUNDARY")
        content = os.read(descriptor, 64 * 1024 + 1)
        after = os.fstat(descriptor)
    except OwnerVoyageOperatorError:
        raise
    except OSError as error:
        raise OwnerVoyageOperatorError("OWNER_VOYAGE_EVIDENCE_BOUNDARY") from error
    finally:
        if descriptor >= 0:
            os.close(descriptor)
    if (
        len(content) != metadata.st_size
        or (after.st_dev, after.st_ino, after.st_size)
        != (metadata.st_dev, metadata.st_ino, metadata.st_size)
    ):
        raise OwnerVoyageOperatorError("OWNER_VOYAGE_EVIDENCE_BOUNDARY")
    return hashlib.sha256(content).hexdigest()


def _parse_instant(value: object) -> datetime:
    if not isinstance(value, str) or not value.endswith("Z"):
        raise ValueError("instant")
    parsed = datetime.fromisoformat(value[:-1] + "+00:00")
    if parsed.tzinfo is None:
        raise ValueError("instant")
    return parsed.astimezone(UTC)


def _format_instant(value: datetime) -> str:
    return value.astimezone(UTC).isoformat(timespec="seconds").replace("+00:00", "Z")


def _reject_duplicate_keys(items: list[tuple[str, object]]) -> dict[str, object]:
    result: dict[str, object] = {}
    for key, value in items:
        if key in result:
            raise ValueError("duplicate JSON key")
        result[key] = value
    return result


def _canonical_json(value: Mapping[str, object]) -> bytes:
    return json.dumps(value, ensure_ascii=False, separators=(",", ":"), sort_keys=True).encode("utf-8")


def _emit_failure(code: str) -> int:
    _emit({"code": code, "state": "FAILED"})
    return 2


def _emit(value: Mapping[str, object]) -> None:
    print(json.dumps(value, ensure_ascii=False, separators=(",", ":"), sort_keys=True))


if __name__ == "__main__":
    raise SystemExit(main())
